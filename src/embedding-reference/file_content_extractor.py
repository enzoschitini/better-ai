import os
import json
from io import BytesIO
from bs4 import BeautifulSoup
from docx import Document
from pptx import Presentation
import fitz  # PyMuPDF      
import pandas as pd
from base64 import b64decode
import yaml
from functools import lru_cache

# Vertex AI imports
import vertexai
from google.oauth2 import service_account
from vertexai.generative_models import GenerativeModel, Part


class FileProcessorStrategy:
    """Interface base para processadores de arquivos ou mídias."""
    def extract_content(self, file_bytes: BytesIO) -> str:
        raise NotImplementedError("O método extract_content deve ser implementado.")


class PDFProcessor(FileProcessorStrategy):
    def extract_content(self, file_bytes: BytesIO) -> str:
        pdf = fitz.open(stream=file_bytes, filetype="pdf")
        content = ''.join([page.get_text() for page in pdf])
        return content


class DocxProcessor(FileProcessorStrategy):
    def extract_content(self, file_bytes: BytesIO) -> str:
        doc = Document(file_bytes)
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells]
                if any(row_text):
                    paragraphs.append(" | ".join(row_text))
        return "\n".join(paragraphs)


class PptxProcessor(FileProcessorStrategy):
    def extract_content(self, file_bytes: BytesIO) -> str:
        prs = Presentation(file_bytes)
        slides_text = []
        for slide in prs.slides:
            slide_lines = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        slide_lines.append(text)
            if slide_lines:
                slides_text.append("\n".join(slide_lines))
        return "\n\n---\n\n".join(slides_text)


class MarkdownProcessor(FileProcessorStrategy):
    def extract_content(self, file_bytes: BytesIO) -> str:
        file_bytes.seek(0)
        return file_bytes.read().decode("utf-8")


class CSVXLSProcessor(FileProcessorStrategy):
    def __init__(self, ext: str):
        self.ext = ext

    def extract_content(self, file_bytes: BytesIO) -> str:
        file_bytes.seek(0)
        if self.ext == "csv":
            df = pd.read_csv(file_bytes)
        elif self.ext == "xlsx":
            df = pd.read_excel(file_bytes, engine='openpyxl')
        elif self.ext == "xls":
            try:
                df = pd.read_excel(file_bytes, engine='xlrd')
            except Exception:
                file_bytes.seek(0)
                df = pd.read_csv(file_bytes)
        elif self.ext == "xml":
            df = pd.read_xml(file_bytes)
        else:
            raise ValueError("Formato tabular não suportado.")
        df = df.fillna('')
        return df.to_markdown(index=False)


class JSONProcessor(FileProcessorStrategy):
    def extract_content(self, file_bytes: BytesIO) -> str:
        file_bytes.seek(0)
        content = json.load(file_bytes)
        return json.dumps(content, indent=2, ensure_ascii=False)


class HTMLProcessor(FileProcessorStrategy):
    def extract_content(self, file_bytes: BytesIO) -> str:
        file_bytes.seek(0)
        html_content = file_bytes.read().decode("utf-8", errors="ignore")
        soup = BeautifulSoup(html_content, "html.parser")
        for tag in soup(["script", "style", "noscript"]):
            tag.decompose()
        text = soup.get_text(separator="\n")
        lines = [line.strip() for line in text.splitlines() if line.strip()]
        return "\n".join(lines)


class TextProcessor(FileProcessorStrategy):
    def extract_content(self, file_bytes: BytesIO) -> str:
        file_bytes.seek(0)
        return file_bytes.read().decode("utf-8")

# Gerenciador de Prompts

class PromptManager:
    """Carrega e fornece prompts de forma centralizada."""

    PROMPT_PATH = "src/knowledge_base/prompts.yaml"

    @staticmethod
    @lru_cache(maxsize=1)
    def load_prompts() -> dict:
        """Carrega o arquivo YAML apenas uma vez."""
        if not os.path.exists(PromptManager.PROMPT_PATH):
            raise FileNotFoundError(f"Arquivo {PromptManager.PROMPT_PATH} não encontrado.")
        with open(PromptManager.PROMPT_PATH, "r", encoding="utf-8") as f:
            return yaml.safe_load(f)

    @classmethod
    def get_prompt(cls, key: str) -> str:
        """Obtém um prompt específico pelo nome."""
        prompts = cls.load_prompts()
        if key not in prompts:
            raise KeyError(f"Prompt '{key}' não encontrado em {cls.PROMPT_PATH}.")
        return prompts[key]


class VertexAIBaseProcessor(FileProcessorStrategy):
    """Classe base para processadores de mídia usando o Vertex AI e prompts externos."""

    def __init__(self, prompt_key: str, mime_type: str):
        self.prompt_key = prompt_key
        self.mime_type = mime_type
        self.prompt_text = PromptManager.get_prompt(prompt_key)
        self._initialize_vertex_ai()
        self.model = GenerativeModel("gemini-2.5-flash")

    def _initialize_vertex_ai(self):
        credentials_b64 = os.environ.get('GOOGLE_SERVICE_CREDENTIALS')
        if not credentials_b64:
            raise EnvironmentError("Credenciais do Vertex AI não encontradas.")
        credentials = json.loads(b64decode(credentials_b64))
        creds = service_account.Credentials.from_service_account_info(credentials)
        vertexai.init(
            project=credentials.get('project_id'),
            location='us-central1',
            credentials=creds
        )

    def _format_usage_metadata(self, usage_metadata):
        text = str(usage_metadata)

        result = {}
        for line in text.strip().split("\n"):
            key, value = line.split(":")
            result[key.strip()] = int(value.strip())

        return result

    def _generate_content(self, media_bytes: bytes, prompt: str, mime_type: str) -> str:
        try:
            part = Part.from_data(media_bytes, mime_type=mime_type)
            gen = self.model.generate_content([part, prompt])

            response = {
                "file_content": gen.text,
                "usage_metadata": self._format_usage_metadata(gen.usage_metadata)
            }
            
            return response
        except Exception as e:
            raise RuntimeError(f"Erro ao processar mídia: {e}")


class ImageProcessor(VertexAIBaseProcessor):
    def __init__(self):
        super().__init__(prompt_key="image_prompt", mime_type="image/png")

    def extract_content(self, file_bytes: BytesIO) -> str:
        file_bytes.seek(0)

        response = self._generate_content(file_bytes.read(), self.prompt_text, self.mime_type)
        return response


class AudioProcessor(VertexAIBaseProcessor):
    def __init__(self):
        super().__init__(prompt_key="audio_prompt", mime_type="audio/mpeg")

    def extract_content(self, file_bytes: BytesIO) -> str:
        file_bytes.seek(0)
        return self._generate_content(file_bytes.read(), self.prompt_text, self.mime_type)


class FileProcessorFactory:
    """Seleciona a estratégia adequada conforme a extensão do arquivo."""

    @staticmethod
    def get_processor(ext: str) -> FileProcessorStrategy:
        ext = ext.lower().replace(".", "")

        mapping = {
            # Arquivos de texto
            "txt": TextProcessor,
            "md": MarkdownProcessor,
            "markdown": MarkdownProcessor,
            "html": HTMLProcessor,

            # Arquivos de documentos
            "pdf": PDFProcessor,
            "docx": DocxProcessor,
            "pptx": PptxProcessor,

            # Arquivos de dados estruturados
            "csv": lambda: CSVXLSProcessor("csv"),
            "xls": lambda: CSVXLSProcessor("xls"),
            "xlsx": lambda: CSVXLSProcessor("xlsx"),
            "xml": lambda: CSVXLSProcessor("xml"),
            "json": JSONProcessor,

            # Arquivos de imagem
            "app": ImageProcessor,
            "jpg": ImageProcessor,
            "jpeg": ImageProcessor,
            "png": ImageProcessor,
            "bmp": ImageProcessor,
            "gif": ImageProcessor,
            "tiff": ImageProcessor,
            "webp": ImageProcessor,

            # Arquivos de áudio
            "mp3": AudioProcessor,
            "wav": AudioProcessor,
            "ogg": AudioProcessor,
            "flac": AudioProcessor,
            "aac": AudioProcessor,
            "m4a": AudioProcessor
        }


        if ext not in mapping:
            raise ValueError(f"Extensão de arquivo não suportada: {ext}")

        processor = mapping[ext]
        return processor() if callable(processor) else processor

# Classe Principal

class FileContentExtractor:
    """
    Classe principal para extrair conteúdo de arquivos e mídias.
    Retorna um JSON no formato {"file_content": "..."}.
    """

    def __init__(self, file_bytes: BytesIO, file_extension: str):
        self.file_bytes = file_bytes
        self.file_extension = file_extension
        self.processor = FileProcessorFactory.get_processor(file_extension)

    def extract(self) -> dict:
        response = self.processor.extract_content(self.file_bytes)
        return {"response": response}
