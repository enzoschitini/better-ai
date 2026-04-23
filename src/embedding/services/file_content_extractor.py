import os
import json
from io import BytesIO
from bs4 import BeautifulSoup
from docx import Document
from pptx import Presentation
import pymupdf as fitz  
import pandas as pd


class FileProcessorStrategy:
    """Interface for file processing strategies."""
    def extract_content(self, file_bytes: BytesIO) -> str:
        raise NotImplementedError("The method extract_content must be implemented by subclasses.")


class PDFProcessor(FileProcessorStrategy):
    def extract_content(self, file_bytes: BytesIO) -> str:
        pdf = fitz.open(stream=file_bytes, filetype="pdf")
        content = ''.join([page.get_text() for page in pdf])
        return content


class DocxProcessor(FileProcessorStrategy):
    def extract_content(self, file_bytes: BytesIO) -> str:
        doc = Document(file_bytes)
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells]
                if any(row_text):
                    paragraphs.append(" | ".join(row_text))
        return "\n".join(paragraphs)


class PptxProcessor(FileProcessorStrategy):
    def extract_content(self, file_bytes: BytesIO) -> str:
        prs = Presentation(file_bytes)
        slides_text = []
        for slide in prs.slides:
            slide_lines = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        slide_lines.append(text)
            if slide_lines:
                slides_text.append("\n".join(slide_lines))
        return "\n\n---\n\n".join(slides_text)


class MarkdownProcessor(FileProcessorStrategy):
    def extract_content(self, file_bytes: BytesIO) -> str:
        file_bytes.seek(0)
        return file_bytes.read().decode("utf-8")


class CSVXLSProcessor(FileProcessorStrategy):
    def __init__(self, ext: str):
        self.ext = ext.lower()

    def _read_csv_safe(self, file_bytes: BytesIO) -> pd.DataFrame:
        strategies = [
            {"sep": ";", "encoding": "utf-8"},
            {"sep": ";", "encoding": "latin-1"},
            {"sep": ",", "encoding": "utf-8"},
        ]

        for strat in strategies:
            try:
                file_bytes.seek(0)
                return pd.read_csv(
                    file_bytes,
                    engine="python",
                    quotechar='"',
                    on_bad_lines="skip",
                    **strat
                )
            except Exception:
                continue

        raise ValueError("Is not possible to read the CSV file with the provided strategies.")

    def extract_content(self, file_bytes: BytesIO, s3_url: str | None = None) -> str:
        file_bytes.seek(0)

        if self.ext == "csv":
            df = self._read_csv_safe(file_bytes)

        elif self.ext == "xlsx":
            df = pd.read_excel(file_bytes, engine="openpyxl")
        elif self.ext == "xls":
            try:
                df = pd.read_excel(file_bytes, engine="xlrd")
            except Exception:
                file_bytes.seek(0)
                df = self._read_csv_safe(file_bytes)

        elif self.ext == "xml":
            df = pd.read_xml(file_bytes)

        else:
            raise ValueError("Unsupported file extension")

        df = df.fillna("").astype(str)
        return df.to_markdown(index=False)


class JSONProcessor(FileProcessorStrategy):
    def extract_content(self, file_bytes: BytesIO) -> str:
        file_bytes.seek(0)
        content = json.load(file_bytes)
        return json.dumps(content, indent=2, ensure_ascii=False)


class HTMLProcessor(FileProcessorStrategy):
    def extract_content(self, file_bytes: BytesIO) -> str:
        file_bytes.seek(0)
        html_content = file_bytes.read().decode("utf-8", errors="ignore")
        soup = BeautifulSoup(html_content, "html.parser")
        for tag in soup(["script", "style", "noscript"]):
            tag.decompose()
        text = soup.get_text(separator="\n")
        lines = [line.strip() for line in text.splitlines() if line.strip()]
        return "\n".join(lines)


class TextProcessor(FileProcessorStrategy):
    def extract_content(self, file_bytes: BytesIO) -> str:
        file_bytes.seek(0)
        return file_bytes.read().decode("utf-8")


class FileProcessorFactory:
    """Selects the appropriate strategy based on the file extension."""

    @staticmethod
    def get_processor(ext: str) -> FileProcessorStrategy:
        ext = ext.lower().replace(".", "")

        mapping = {
            # Text files
            "txt": TextProcessor,
            "md": MarkdownProcessor,
            "markdown": MarkdownProcessor,
            "html": HTMLProcessor,

            # Document files
            "pdf": PDFProcessor,
            "doc": DocxProcessor,
            "docx": DocxProcessor,
            "ppt": PptxProcessor,
            "pptx": PptxProcessor,

            # Structured data files
            "csv": lambda: CSVXLSProcessor("csv"),
            "xls": lambda: CSVXLSProcessor("xls"),
            "xlsx": lambda: CSVXLSProcessor("xlsx"),
            "xml": lambda: CSVXLSProcessor("xml"),
            "json": JSONProcessor
        }


        if ext not in mapping:
            raise ValueError(f"Unsupported file extension: {ext}")

        processor = mapping[ext]
        return processor() if callable(processor) else processor

# Main Class

class FileContentExtractor:
    """
    Main class for extracting content from files and media.
    Returns a JSON in the format {"file_content": "..."}.
    """

    def __init__(self, file_bytes: BytesIO, file_extension: str):
        self.file_bytes = file_bytes
        self.file_extension = file_extension
        self.processor = FileProcessorFactory.get_processor(file_extension)

    def extract(self) -> dict:
        response = self.processor.extract_content(self.file_bytes)
        return {"file_content": response}

if __name__ == "__main__":
    path = "local/test files"

    # Exemplo de uso
    with open(f"{path}/Candidatura.pdf", "rb") as f:
        file_bytes = BytesIO(f.read())

    extractor = FileContentExtractor(file_bytes, "pdf")
    result = extractor.extract()
    print(result["file_content"][:500])  # Imprime os primeiros 500 caracteres do conteúdo extraído

# python -m src.embedding.services.file_content_extractor